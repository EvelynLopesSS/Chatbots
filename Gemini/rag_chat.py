import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
import os
from pptx import Presentation

st.set_page_config(page_title="Document Genie", layout="wide")

st.markdown("""
## Document Genie: Get instant insights from your Documents

This chatbot is built using the Retrieval-Augmented Generation (RAG) framework, leveraging Google's Generative AI model Gemini-PRO. It processes uploaded PDF documents by breaking them down into manageable chunks, creates a searchable vector store, and generates accurate answers to user queries. This advanced approach ensures high-quality, contextually relevant responses for an efficient and effective user experience.

### How It Works

Follow these simple steps to interact with the chatbot:

1. **Enter Your API Key**: You'll need a Google API key for the chatbot to access Google's Generative AI models. Obtain your API key https://makersuite.google.com/app/apikey.

2. **Upload Your Documents**: The system accepts multiple PDF files at once, analyzing the content to provide comprehensive insights.

3. **Ask a Question**: After processing the documents, ask any question related to the content of your uploaded documents for a precise answer.
""")



# This is the first API key input; no need to repeat it in the main function.
api_key = st.sidebar.text_input("Enter your Google API Key:", type="password", key="api_key_input")

def extract_text_from_PDF(file_path):
    text = ""
    for pdf in file_path:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_PPTX(file_path):
    texto = ""
    try:
        for file in file_path:
            apresentacao = Presentation(file.getvalue())
            for slide in apresentacao.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
    except Exception as e:
        print("Erro ao extrair texto do arquivo PPTX:", e)
    return texto

def extract_text_from_Files(file_path):
    for file in file_path:
        ext = os.path.splitext(file.name)[1].lower()
        if ext == ".pdf":
            raw_text = extract_text_from_PDF([file])
        elif ext == ".pptx":
            raw_text = extract_text_from_PPTX([file])
        else:
            st.error(f"Unsupported file format: {ext}")
            return
    st.success("Done")

def process_uploaded_files(uploaded_files):

    if not os.path.exists('files'):
        os.mkdir('files')

    for uploaded_file in uploaded_files:
        if uploaded_file is not None:
            if not os.path.isfile("files/" + uploaded_file.name):
                with st.status("Processing your document..."):
                    bytes_data = uploaded_file.read()
                    with open("files/" + uploaded_file.name, "wb") as f:
                        f.write(bytes_data)
                extract_text_from_Files("files/" + uploaded_file.name)
                st.success("Your document has been processed.")
                st.balloons()     

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks, api_key):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=api_key)
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3, google_api_key=api_key)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question, api_key):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=api_key)
    new_db = FAISS.load_local("faiss_index", embeddings)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    st.write("Reply: ", response["output_text"])

def main():
    st.header("AI clone chatbot💁")

    user_question = st.text_input("Ask a Question from the PDF Files", key="user_question")

    if user_question and api_key:  # Ensure API key and user question are provided
        user_input(user_question, api_key)

    with st.sidebar:
        st.title("Menu:")
        uploaded_files = st.file_uploader("Upload your PDF Files and Click on the Submit & Process Button", accept_multiple_files=True, key="pdf_uploader")
        if st.button("Submit & Process", key="process_button") and api_key:  # Check if API key is provided before processing
            with st.spinner("Processing..."):
                process_uploaded_files(uploaded_files)
                raw_text = extract_text_from_Files(["files/" + uploaded_file.name for uploaded_file in uploaded_files])
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks, api_key)
                st.success("Done")

if __name__ == "__main__":
    main()